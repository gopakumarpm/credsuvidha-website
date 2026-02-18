"""
CredSuvidha Brand Kit — PowerPoint Presentation Generator
Generates a professional brand guidelines PPTX using python-pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from datetime import datetime

# ─── PATHS ───
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(os.path.dirname(BASE_DIR))
LOGO_PATH = os.path.join(PROJECT_DIR, "logo.png")
OUTPUT_PATH = os.path.join(BASE_DIR, "CredSuvidha-BrandKit.pptx")

# ─── CREDSUVIDHA BRAND COLORS ───
NAVY_DARK    = RGBColor(0x14, 0x28, 0x57)
NAVY         = RGBColor(0x19, 0x3F, 0x8F)
BLUE_PRIMARY = RGBColor(0x1A, 0x6E, 0xF5)
BLUE_BRIGHT  = RGBColor(0x33, 0x8D, 0xFF)
BLUE_LIGHT   = RGBColor(0x59, 0xB0, 0xFF)
BLUE_PALE    = RGBColor(0xBC, 0xE0, 0xFF)
BLUE_TINT    = RGBColor(0xEE, 0xF7, 0xFF)

ACCENT_PRI   = RGBColor(0xF9, 0x73, 0x16)
ACCENT_LT    = RGBColor(0xFB, 0x92, 0x3C)

EMERALD      = RGBColor(0x10, 0xB9, 0x81)
EMERALD_DK   = RGBColor(0x05, 0x96, 0x69)

LOGO_NAVY    = RGBColor(0x1B, 0x3A, 0x5C)
LOGO_GOLD    = RGBColor(0xC5, 0x96, 0x1E)
LOGO_BG      = RGBColor(0xFA, 0xF6, 0xF1)

TEXT_DARK    = RGBColor(0x1E, 0x29, 0x3B)
TEXT_GRAY    = RGBColor(0x64, 0x74, 0x8B)
TEXT_MUTED   = RGBColor(0x94, 0xA3, 0xB8)
BORDER       = RGBColor(0xE2, 0xE8, 0xF0)
BG_ALT       = RGBColor(0xF8, 0xFA, 0xFC)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x00, 0x00, 0x00)


def set_font(run, size_pt, bold=False, italic=False, color=TEXT_DARK, name='Calibri'):
    """Set font properties on a run."""
    run.font.name = name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_shape_with_fill(slide, left, top, width, height, fill_color, shape_type=MSO_SHAPE.RECTANGLE):
    """Add a filled shape to a slide."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=16, bold=False,
                italic=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, font_size, bold, italic, color, font_name)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=TEXT_DARK,
                    bullet_color=None, spacing=Pt(6)):
    """Add a bulleted list to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        # Add bullet character
        bullet_run = p.add_run()
        bullet_run.text = "●  "
        set_font(bullet_run, font_size - 2, color=bullet_color or BLUE_PRIMARY)

        run = p.add_run()
        run.text = item
        set_font(run, font_size, color=color)

    return txBox


def add_accent_bar(slide, top=Inches(0), color=BLUE_PRIMARY, height=Pt(4)):
    """Add a thin accent bar across the top of a slide."""
    add_shape_with_fill(slide, Inches(0), top, Inches(13.333), height, color)


def add_footer_bar(slide):
    """Add CredSuvidha branded footer bar."""
    bar = add_shape_with_fill(slide, Inches(0), Inches(7.08), Inches(13.333), Inches(0.42), NAVY_DARK)

    # Left text
    add_textbox(slide, Inches(0.4), Inches(7.13), Inches(5), Inches(0.3),
                "CredSuvidha — Trusted Partner. Swift Solutions.",
                font_size=9, color=BLUE_LIGHT, italic=True)

    # Right text
    add_textbox(slide, Inches(9.5), Inches(7.13), Inches(3.5), Inches(0.3),
                "www.credsuvidha.com",
                font_size=9, color=LOGO_GOLD, alignment=PP_ALIGN.RIGHT)


def add_slide_number(slide, number, total=12):
    """Add slide number."""
    add_textbox(slide, Inches(6.2), Inches(7.15), Inches(1), Inches(0.25),
                f"{number}/{total}", font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)


def add_logo_header(slide):
    """Add small logo in top-right corner of content slides."""
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), width=Inches(2.2))


# ═══════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════

def build_title_slide(prs):
    """Slide 1: Title / Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY_DARK

    # Gold accent stripe at bottom of top section
    add_shape_with_fill(slide, Inches(0), Inches(4.5), Inches(6.667), Pt(5), BLUE_PRIMARY)
    add_shape_with_fill(slide, Inches(6.667), Inches(4.5), Inches(6.666), Pt(5), LOGO_GOLD)

    # Logo (centered, large)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(3.8), Inches(0.6), width=Inches(5.5))

    # Title text
    add_textbox(slide, Inches(1), Inches(4.7), Inches(11.333), Inches(0.8),
                "Brand Kit", font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(1), Inches(5.4), Inches(11.333), Inches(0.6),
                "Visual Identity & Brand Guidelines", font_size=20, color=BLUE_LIGHT, alignment=PP_ALIGN.CENTER)

    # Version & date
    add_textbox(slide, Inches(1), Inches(6.1), Inches(11.333), Inches(0.4),
                f"Version 1.0  •  {datetime.now().strftime('%B %Y')}  •  www.credsuvidha.com",
                font_size=11, italic=True, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # Bottom accent
    add_shape_with_fill(slide, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3), BLUE_PRIMARY)

    # Speaker notes
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("[KEY POINT]: Welcome to the CredSuvidha Brand Kit presentation.\n"
                  "[DATA]: This guide covers our visual identity including logo, colors, typography, "
                  "and UI components.\n"
                  "[TRANSITION]: Let's start with an overview of our brand.")


def build_brand_overview(prs):
    """Slide 2: Brand Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_ALT

    add_accent_bar(slide)
    add_logo_header(slide)
    add_footer_bar(slide)
    add_slide_number(slide, 2)

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                "Brand Overview", font_size=32, bold=True, color=NAVY_DARK)

    # Mission statement box
    mission_box = add_shape_with_fill(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.0), BLUE_TINT)
    mission_box.line.color.rgb = BLUE_PRIMARY
    mission_box.line.width = Pt(1)

    add_textbox(slide, Inches(1.2), Inches(1.5), Inches(11), Inches(0.8),
                '"Smart Financial Solutions for a Better Tomorrow"',
                font_size=22, bold=True, color=NAVY_DARK, alignment=PP_ALIGN.CENTER)

    # Tagline
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.5))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = "Trusted Partner. "; set_font(r1, 18, bold=True, color=NAVY_DARK)
    r2 = p.add_run(); r2.text = "Swift Solutions."; set_font(r2, 18, bold=True, color=LOGO_GOLD)

    # Key Stats — 4 cards
    stats = [
        ("500+ Cr", "Loans Disbursed", BLUE_PRIMARY),
        ("10,000+", "Happy Customers", LOGO_GOLD),
        ("50+", "Banking Partners", EMERALD),
        ("24/7", "Expert Support", ACCENT_PRI),
    ]

    card_width = Inches(2.6)
    card_height = Inches(1.5)
    start_x = Inches(0.8)
    spacing = Inches(0.2)
    card_y = Inches(3.3)

    for i, (number, label, color) in enumerate(stats):
        x = start_x + i * (card_width + spacing)

        # Card background
        card = add_shape_with_fill(slide, x, card_y, card_width, card_height, WHITE)
        card.line.color.rgb = BORDER
        card.line.width = Pt(0.5)

        # Top accent
        add_shape_with_fill(slide, x, card_y, card_width, Pt(4), color)

        # Number
        add_textbox(slide, x, card_y + Inches(0.2), card_width, Inches(0.6),
                    number, font_size=28, bold=True, color=color, alignment=PP_ALIGN.CENTER)

        # Label
        add_textbox(slide, x, card_y + Inches(0.85), card_width, Inches(0.4),
                    label, font_size=12, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # Description
    add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.0),
                "CredSuvidha is a modern financial services facilitator providing smart solutions across "
                "loans, credit cards, and insurance. We partner with 50+ banks to offer the best rates "
                "and fastest approvals for every Indian customer.",
                font_size=13, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # Industry badge
    badge = add_shape_with_fill(slide, Inches(5.0), Inches(6.3), Inches(3.3), Inches(0.35), BLUE_PRIMARY,
                                MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(5.0), Inches(6.32), Inches(3.3), Inches(0.3),
                "Financial Services (Fintech)", font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    notes = slide.notes_slide.notes_text_frame
    notes.text = ("[KEY POINT]: CredSuvidha is a fintech facilitator providing loans, credit cards, and insurance.\n"
                  "[DATA]: 500+ Cr disbursed, 10K+ customers, 50+ banking partners, 24/7 support.\n"
                  "[TRANSITION]: Let's look at our visual identity starting with the logo.")


def build_logo_slide(prs):
    """Slide 3: Logo Guidelines"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_ALT

    add_accent_bar(slide)
    add_logo_header(slide)
    add_footer_bar(slide)
    add_slide_number(slide, 3)

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                "Logo", font_size=32, bold=True, color=NAVY_DARK)

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.5),
                "Interlocking C and S letterforms with an upward growth arrow — symbolizing financial progress and partnership.",
                font_size=13, color=TEXT_GRAY)

    # Logo on white background
    if os.path.exists(LOGO_PATH):
        # White card
        card1 = add_shape_with_fill(slide, Inches(0.8), Inches(1.8), Inches(5.4), Inches(2.8), WHITE)
        card1.line.color.rgb = BORDER; card1.line.width = Pt(0.5)
        slide.shapes.add_picture(LOGO_PATH, Inches(1.5), Inches(2.0), width=Inches(4.0))
        add_textbox(slide, Inches(0.8), Inches(4.3), Inches(5.4), Inches(0.3),
                    "Primary — Light Background", font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

        # Cream card
        card2 = add_shape_with_fill(slide, Inches(6.5), Inches(1.8), Inches(5.9), Inches(2.8), LOGO_BG)
        card2.line.color.rgb = BORDER; card2.line.width = Pt(0.5)
        slide.shapes.add_picture(LOGO_PATH, Inches(7.4), Inches(2.0), width=Inches(4.0))
        add_textbox(slide, Inches(6.5), Inches(4.3), Inches(5.9), Inches(0.3),
                    "On Brand Cream — #FAF6F1", font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # Logo elements table area
    add_textbox(slide, Inches(0.8), Inches(4.85), Inches(4), Inches(0.4),
                "Logo Elements", font_size=16, bold=True, color=NAVY_DARK)

    elements = [
        ("Symbol", "Interlocking C (navy) & S (gold) with upward arrow"),
        ("Wordmark", '"CREDSUVIDHA.COM" in navy blue uppercase'),
        ("Tagline", '"TRUSTED PARTNER. SWIFT SOLUTIONS." in gold'),
        ("Min Size", "120px wide (digital) / 30mm (print)"),
    ]

    for i, (label, desc) in enumerate(elements):
        y = Inches(5.25) + i * Inches(0.35)
        add_textbox(slide, Inches(1.0), y, Inches(1.5), Inches(0.3),
                    f"● {label}:", font_size=11, bold=True, color=BLUE_PRIMARY)
        add_textbox(slide, Inches(2.6), y, Inches(4.5), Inches(0.3),
                    desc, font_size=11, color=TEXT_DARK)

    # Do's and Don'ts
    add_textbox(slide, Inches(7.5), Inches(4.85), Inches(2.5), Inches(0.4),
                "✓  Do", font_size=16, bold=True, color=EMERALD_DK)

    dos = ["Use on white, cream, or light backgrounds",
           "Use inverted version on dark backgrounds",
           "Maintain proportions when scaling"]
    for i, item in enumerate(dos):
        add_textbox(slide, Inches(7.7), Inches(5.3) + i * Inches(0.32), Inches(5), Inches(0.3),
                    f"●  {item}", font_size=10, color=TEXT_DARK)

    add_textbox(slide, Inches(10.2), Inches(4.85), Inches(2.5), Inches(0.4),
                "✗  Don't", font_size=16, bold=True, color=RGBColor(0xDC, 0x26, 0x26))

    donts = ["Stretch or distort the logo",
             "Change the logo colors",
             "Place on busy backgrounds"]
    for i, item in enumerate(donts):
        add_textbox(slide, Inches(10.4), Inches(5.3) + i * Inches(0.32), Inches(3), Inches(0.3),
                    f"●  {item}", font_size=10, color=TEXT_DARK)

    notes = slide.notes_slide.notes_text_frame
    notes.text = ("[KEY POINT]: The logo features interlocking C and S with a growth arrow.\n"
                  "[DATA]: Navy blue = trust/stability, Gold = prosperity/value.\n"
                  "[TRANSITION]: Now let's examine our color palette in detail.")


def build_color_palette_slide(prs):
    """Slide 4: Color Palette — Brand Blue"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_ALT

    add_accent_bar(slide)
    add_logo_header(slide)
    add_footer_bar(slide)
    add_slide_number(slide, 4)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                "Color Palette", font_size=32, bold=True, color=NAVY_DARK)

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4),
                "Built around trust (navy blue), warmth (gold/amber), energy (orange), and growth (emerald green).",
                font_size=13, color=TEXT_GRAY)

    # ── Brand Blue row ──
    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(3), Inches(0.3),
                "Primary — Brand Blue", font_size=14, bold=True, color=NAVY_DARK)

    brand_blues = [
        ("#142857", "950"), ("#193f8f", "900"), ("#1747b6", "800"), ("#1458e1", "700"),
        ("#1a6ef5", "600★"), ("#338dff", "500"), ("#59b0ff", "400"),
        ("#8ecdff", "300"), ("#bce0ff", "200"), ("#eef7ff", "50"),
    ]

    sw = Inches(1.15)
    sh = Inches(0.65)
    sx = Inches(0.8)
    sy = Inches(1.95)

    for i, (hex_code, label) in enumerate(brand_blues):
        x = sx + i * (sw + Inches(0.08))
        r, g, b = int(hex_code[1:3], 16), int(hex_code[3:5], 16), int(hex_code[5:7], 16)
        swatch = add_shape_with_fill(slide, x, sy, sw, sh, RGBColor(r, g, b))
        swatch.line.color.rgb = BORDER; swatch.line.width = Pt(0.5)

        text_col = WHITE if (r + g + b) < 400 else TEXT_DARK
        add_textbox(slide, x, sy + sh + Pt(2), sw, Inches(0.2),
                    label, font_size=8, bold=True, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, sy + sh + Inches(0.18), sw, Inches(0.2),
                    hex_code.upper(), font_size=7, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # ── Logo Colors ──
    add_textbox(slide, Inches(0.8), Inches(3.1), Inches(4), Inches(0.3),
                "Logo Colors", font_size=14, bold=True, color=NAVY_DARK)

    logo_colors = [("#1B3A5C", "Logo Navy"), ("#C5961E", "Logo Gold"), ("#FAF6F1", "Logo BG")]
    for i, (hex_code, label) in enumerate(logo_colors):
        x = Inches(0.8) + i * Inches(2.5)
        r, g, b = int(hex_code[1:3], 16), int(hex_code[3:5], 16), int(hex_code[5:7], 16)
        swatch = add_shape_with_fill(slide, x, Inches(3.4), Inches(2.2), Inches(0.55), RGBColor(r, g, b))
        swatch.line.color.rgb = BORDER; swatch.line.width = Pt(0.5)
        add_textbox(slide, x, Inches(3.98), Inches(2.2), Inches(0.3),
                    f"{label}  {hex_code}", font_size=9, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Accent Orange ──
    add_textbox(slide, Inches(0.8), Inches(4.45), Inches(4), Inches(0.3),
                "Accent — Orange", font_size=14, bold=True, color=NAVY_DARK)

    accent_colors = [
        ("#c2410c", "700"), ("#ea580c", "600"), ("#f97316", "500★"),
        ("#fb923c", "400"), ("#fdba74", "300"), ("#fff7ed", "50"),
    ]
    for i, (hex_code, label) in enumerate(accent_colors):
        x = Inches(0.8) + i * (Inches(2.0) + Inches(0.08))
        r, g, b = int(hex_code[1:3], 16), int(hex_code[3:5], 16), int(hex_code[5:7], 16)
        swatch = add_shape_with_fill(slide, x, Inches(4.8), Inches(2.0), Inches(0.55), RGBColor(r, g, b))
        swatch.line.color.rgb = BORDER; swatch.line.width = Pt(0.5)
        text_col = WHITE if (r + g + b) < 400 else TEXT_DARK
        add_textbox(slide, x, Inches(5.38), Inches(2.0), Inches(0.25),
                    f"{label}  {hex_code.upper()}", font_size=8, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Emerald ──
    add_textbox(slide, Inches(0.8), Inches(5.8), Inches(4), Inches(0.3),
                "Success — Emerald Green", font_size=14, bold=True, color=NAVY_DARK)

    emerald_colors = [
        ("#059669", "600"), ("#10b981", "500"), ("#34d399", "400"), ("#ecfdf5", "50"),
    ]
    for i, (hex_code, label) in enumerate(emerald_colors):
        x = Inches(0.8) + i * Inches(2.5)
        r, g, b = int(hex_code[1:3], 16), int(hex_code[3:5], 16), int(hex_code[5:7], 16)
        swatch = add_shape_with_fill(slide, x, Inches(6.1), Inches(2.2), Inches(0.5), RGBColor(r, g, b))
        swatch.line.color.rgb = BORDER; swatch.line.width = Pt(0.5)
        add_textbox(slide, x, Inches(6.63), Inches(2.2), Inches(0.25),
                    f"{label}  {hex_code.upper()}", font_size=8, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    notes = slide.notes_slide.notes_text_frame
    notes.text = ("[KEY POINT]: Our color system has 4 families — brand blue, logo colors, accent orange, emerald.\n"
                  "[DATA]: Primary brand color is #1A6EF5 (Brand 600). Logo uses #1B3A5C navy and #C5961E gold.\n"
                  "[TRANSITION]: Next, let's look at our gradient specifications.")


def build_gradients_slide(prs):
    """Slide 5: Gradients"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_ALT

    add_accent_bar(slide)
    add_logo_header(slide)
    add_footer_bar(slide)
    add_slide_number(slide, 5)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                "Gradients", font_size=32, bold=True, color=NAVY_DARK)

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4),
                "Signature gradients used across buttons, backgrounds, and accent elements.",
                font_size=13, color=TEXT_GRAY)

    gradients = [
        ("Primary CTA", "linear-gradient(135deg, #1a6ef5, #338dff)",
         [(BLUE_PRIMARY, 0.5), (BLUE_BRIGHT, 0.5)]),
        ("Hero / Header", "linear-gradient(135deg, #142857 → #1a6ef5 → #338dff)",
         [(NAVY_DARK, 0.33), (BLUE_PRIMARY, 0.34), (BLUE_BRIGHT, 0.33)]),
        ("Accent CTA", "linear-gradient(135deg, #f97316, #fb923c)",
         [(ACCENT_PRI, 0.5), (ACCENT_LT, 0.5)]),
        ("Gradient Text", "linear-gradient(135deg, #1a6ef5 → #338dff → #f97316)",
         [(BLUE_PRIMARY, 0.33), (BLUE_BRIGHT, 0.34), (ACCENT_PRI, 0.33)]),
        ("Success", "linear-gradient(135deg, #059669, #10b981)",
         [(EMERALD_DK, 0.5), (EMERALD, 0.5)]),
        ("Logo Badge", "linear-gradient(to bottom right, #1a6ef5, #59b0ff)",
         [(BLUE_PRIMARY, 0.5), (BLUE_LIGHT, 0.5)]),
    ]

    card_w = Inches(5.8)
    card_h = Inches(0.7)
    start_y = Inches(1.65)

    for i, (name, css, color_stops) in enumerate(gradients):
        col = 0 if i < 3 else 1
        row = i % 3
        x = Inches(0.8) + col * (card_w + Inches(0.3))
        y = start_y + row * (card_h + Inches(0.55))

        # Simulate gradient with segments
        total_w = card_w
        cx = x
        for color, fraction in color_stops:
            seg_w = int(total_w * fraction)
            seg = add_shape_with_fill(slide, cx, y, seg_w, card_h, color)
            seg.line.fill.background()
            cx += seg_w

        # Label below
        add_textbox(slide, x, y + card_h + Pt(2), card_w, Inches(0.2),
                    f"{name}", font_size=11, bold=True, color=NAVY_DARK)
        add_textbox(slide, x, y + card_h + Inches(0.22), card_w, Inches(0.2),
                    css, font_size=8, color=TEXT_MUTED)

    notes = slide.notes_slide.notes_text_frame
    notes.text = ("[KEY POINT]: We have 6 signature gradients for different use cases.\n"
                  "[DATA]: Primary CTA uses brand-600 → brand-500, Hero uses navy → brand-600 → brand-500.\n"
                  "[TRANSITION]: Let's examine our typography system.")


def build_typography_slide(prs):
    """Slide 6: Typography"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_ALT

    add_accent_bar(slide)
    add_logo_header(slide)
    add_footer_bar(slide)
    add_slide_number(slide, 6)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                "Typography", font_size=32, bold=True, color=NAVY_DARK)

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4),
                "Dual typeface system — clean sans-serif for body, elegant serif for display.",
                font_size=13, color=TEXT_GRAY)

    # ── Inter section ──
    inter_box = add_shape_with_fill(slide, Inches(0.8), Inches(1.65), Inches(7.5), Inches(3.2), WHITE)
    inter_box.line.color.rgb = BORDER; inter_box.line.width = Pt(0.5)

    add_textbox(slide, Inches(1.0), Inches(1.75), Inches(7), Inches(0.4),
                "Inter — Primary Typeface (Sans-Serif)", font_size=18, bold=True, color=BLUE_PRIMARY)

    add_textbox(slide, Inches(1.0), Inches(2.15), Inches(7), Inches(0.3),
                "Body text, navigation, buttons, labels, all UI elements",
                font_size=11, italic=True, color=TEXT_GRAY)

    weights = [
        ("Light 300", False), ("Regular 400", False), ("Medium 500", False),
        ("SemiBold 600", True), ("Bold 700", True), ("ExtraBold 800", True),
    ]
    for i, (weight, bold) in enumerate(weights):
        y = Inches(2.55) + i * Inches(0.35)
        add_textbox(slide, Inches(1.2), y, Inches(1.8), Inches(0.3),
                    weight, font_size=10, color=TEXT_MUTED)
        add_textbox(slide, Inches(3.0), y, Inches(5), Inches(0.3),
                    "Smart Financial Solutions for a Better Tomorrow",
                    font_size=13, bold=bold, color=TEXT_DARK)

    # ── Playfair Display ──
    playfair_box = add_shape_with_fill(slide, Inches(0.8), Inches(5.0), Inches(7.5), Inches(1.3), WHITE)
    playfair_box.line.color.rgb = BORDER; playfair_box.line.width = Pt(0.5)

    add_textbox(slide, Inches(1.0), Inches(5.1), Inches(7), Inches(0.4),
                "Playfair Display — Display Typeface (Serif)", font_size=18, bold=True, color=BLUE_PRIMARY)

    add_textbox(slide, Inches(1.0), Inches(5.5), Inches(7), Inches(0.3),
                "Hero headlines and special emphasis. Bold (700) and ExtraBold (800).",
                font_size=11, italic=True, color=TEXT_GRAY)

    add_textbox(slide, Inches(1.0), Inches(5.85), Inches(7), Inches(0.35),
                "Trusted Partner. Swift Solutions.",
                font_size=20, bold=True, color=NAVY_DARK, font_name='Calibri')

    # ── Type Scale (right side) ──
    add_textbox(slide, Inches(8.6), Inches(1.65), Inches(4), Inches(0.4),
                "Type Scale", font_size=18, bold=True, color=NAVY_DARK)

    scale_items = [
        ("Hero H1", "48-60px", "ExtraBold"),
        ("Section H2", "36-40px", "Bold"),
        ("Card H3", "20-24px", "SemiBold"),
        ("Body", "14-16px", "Regular"),
        ("Small/Label", "12-13px", "Medium"),
        ("Caption", "10-11px", "Medium, uppercase"),
    ]

    scale_box = add_shape_with_fill(slide, Inches(8.6), Inches(2.1), Inches(4.2), Inches(3.4), WHITE)
    scale_box.line.color.rgb = BORDER; scale_box.line.width = Pt(0.5)

    for i, (element, size, weight) in enumerate(scale_items):
        y = Inches(2.2) + i * Inches(0.52)
        add_textbox(slide, Inches(8.8), y, Inches(1.5), Inches(0.25),
                    element, font_size=11, bold=True, color=BLUE_PRIMARY)
        add_textbox(slide, Inches(10.3), y, Inches(1.0), Inches(0.25),
                    size, font_size=10, color=TEXT_DARK)
        add_textbox(slide, Inches(11.3), y, Inches(1.3), Inches(0.25),
                    weight, font_size=10, color=TEXT_GRAY)

    # Google Fonts URL
    add_textbox(slide, Inches(8.6), Inches(5.6), Inches(4.2), Inches(0.8),
                "Google Fonts:\nInter: wght@300;400;500;600;700;800;900\nPlayfair Display: wght@700;800",
                font_size=9, color=TEXT_MUTED)

    notes = slide.notes_slide.notes_text_frame
    notes.text = ("[KEY POINT]: Two fonts — Inter for everything, Playfair Display for special headlines.\n"
                  "[DATA]: Inter has 7 weights (300-900). Playfair Display has Bold and ExtraBold.\n"
                  "[TRANSITION]: Let's look at how these translate into UI components.")


def build_buttons_ui_slide(prs):
    """Slide 7: Buttons & UI Components"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_ALT

    add_accent_bar(slide)
    add_logo_header(slide)
    add_footer_bar(slide)
    add_slide_number(slide, 7)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                "Buttons & UI Components", font_size=32, bold=True, color=NAVY_DARK)

    # Button variants
    add_textbox(slide, Inches(0.8), Inches(1.3), Inches(4), Inches(0.4),
                "Button Variants", font_size=18, bold=True, color=NAVY_DARK)

    buttons = [
        ("Get Started →", BLUE_PRIMARY, WHITE, "Primary"),
        ("Explore Services", WHITE, BLUE_PRIMARY, "Secondary"),
        ("Apply Now", ACCENT_PRI, WHITE, "Accent"),
        ("Learn More", NAVY_DARK, WHITE, "Dark"),
    ]

    for i, (text, bg_color, text_color, label) in enumerate(buttons):
        x = Inches(0.8) + i * Inches(2.8)
        y = Inches(1.8)

        btn = add_shape_with_fill(slide, x, y, Inches(2.4), Inches(0.55), bg_color,
                                  MSO_SHAPE.ROUNDED_RECTANGLE)
        if bg_color == WHITE:
            btn.line.color.rgb = BLUE_PRIMARY
            btn.line.width = Pt(2)

        add_textbox(slide, x, y + Pt(4), Inches(2.4), Inches(0.45),
                    text, font_size=13, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.6), Inches(2.4), Inches(0.2),
                    label, font_size=9, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # Specs
    add_textbox(slide, Inches(0.8), Inches(2.7), Inches(11), Inches(0.3),
                "Border-radius: 100px (pill)  •  Padding: 12px 28px  •  Font: Inter SemiBold 14px  •  Transition: 300ms",
                font_size=10, italic=True, color=TEXT_MUTED)

    # UI Effects section
    add_textbox(slide, Inches(0.8), Inches(3.2), Inches(4), Inches(0.4),
                "UI Effects", font_size=18, bold=True, color=NAVY_DARK)

    effects = [
        ("Glass Morphism", "bg: rgba(255,255,255,0.08)\nbackdrop-filter: blur(20px)\nborder: 1px solid rgba(255,255,255,0.15)", "Hero stats, overlaid cards"),
        ("Card Hover", "transform: translateY(-8px)\nbox-shadow: 0 25px 60px rgba(0,0,0,0.12)", "Service cards, feature cards"),
        ("Pulse Glow", "box-shadow: 0 0 20-40px\nrgba(26,110,245, 0.3-0.6)", "Active states, attention draw"),
        ("Scroll Reveal", "opacity: 0→1\ntranslateY(30px)→0\ntransition: 0.8s ease", "Section entrance animations"),
    ]

    for i, (title, css, usage) in enumerate(effects):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.3)
        y = Inches(3.7) + row * Inches(1.55)

        card = add_shape_with_fill(slide, x, y, Inches(5.9), Inches(1.35), WHITE)
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5)

        # Top accent
        add_shape_with_fill(slide, x, y, Inches(5.9), Pt(3), BLUE_PRIMARY)

        add_textbox(slide, x + Inches(0.2), y + Inches(0.1), Inches(2), Inches(0.3),
                    title, font_size=13, bold=True, color=NAVY_DARK)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.4), Inches(3.2), Inches(0.8),
                    css, font_size=9, color=TEXT_GRAY)
        add_textbox(slide, x + Inches(3.5), y + Inches(0.4), Inches(2.2), Inches(0.8),
                    f"Usage:\n{usage}", font_size=10, color=BLUE_PRIMARY)

    notes = slide.notes_slide.notes_text_frame
    notes.text = ("[KEY POINT]: 4 button variants and 4 key UI effects define our component system.\n"
                  "[DATA]: All buttons use pill shape (100px radius), transitions at 300ms.\n"
                  "[TRANSITION]: Let's discuss our core brand values.")


def build_brand_values_slide(prs):
    """Slide 8: Brand Values"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_ALT

    add_accent_bar(slide)
    add_logo_header(slide)
    add_footer_bar(slide)
    add_slide_number(slide, 8)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                "Brand Values", font_size=32, bold=True, color=NAVY_DARK)

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4),
                "Four pillars that define everything we do at CredSuvidha.",
                font_size=13, color=TEXT_GRAY)

    values = [
        ("Trust", "RBI regulated partners, IRDAI registered, ISO 27001 compliant. Security and compliance are at our core.",
         BLUE_PRIMARY, "🛡"),
        ("Speed", "Swift paperless loan approvals. Quick turnaround powered by 50+ banking partners.",
         LOGO_GOLD, "⚡"),
        ("Expertise", "24/7 expert support. 500+ Cr loans disbursed. 10,000+ happy customers served.",
         EMERALD, "🎯"),
        ("Simplicity", "Clean, modern, and accessible. Making financial decisions easy for every Indian.",
         ACCENT_PRI, "✨"),
    ]

    card_w = Inches(5.8)
    card_h = Inches(2.3)

    for i, (title, desc, color, icon) in enumerate(values):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * (card_w + Inches(0.3))
        y = Inches(1.65) + row * (card_h + Inches(0.25))

        # Card
        card = add_shape_with_fill(slide, x, y, card_w, card_h, WHITE)
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5)

        # Top color accent
        add_shape_with_fill(slide, x, y, card_w, Pt(5), color)

        # Icon circle
        circle = add_shape_with_fill(slide, x + Inches(0.3), y + Inches(0.3),
                                     Inches(0.7), Inches(0.7), color, MSO_SHAPE.OVAL)

        # Title
        add_textbox(slide, x + Inches(1.2), y + Inches(0.35), Inches(4), Inches(0.4),
                    title, font_size=22, bold=True, color=NAVY_DARK)

        # Description
        add_textbox(slide, x + Inches(0.3), y + Inches(1.15), card_w - Inches(0.6), Inches(0.9),
                    desc, font_size=13, color=TEXT_GRAY)

    notes = slide.notes_slide.notes_text_frame
    notes.text = ("[KEY POINT]: Trust, Speed, Expertise, and Simplicity are our four brand pillars.\n"
                  "[DATA]: These values guide every decision from product design to customer service.\n"
                  "[TRANSITION]: Let's wrap up with contact info and compliance details.")


def build_contact_compliance_slide(prs):
    """Slide 9: Contact & Compliance"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_ALT

    add_accent_bar(slide)
    add_logo_header(slide)
    add_footer_bar(slide)
    add_slide_number(slide, 9)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                "Contact & Compliance", font_size=32, bold=True, color=NAVY_DARK)

    # Contact Card
    contact_card = add_shape_with_fill(slide, Inches(0.8), Inches(1.4), Inches(5.8), Inches(3.5), WHITE)
    contact_card.line.color.rgb = BORDER; contact_card.line.width = Pt(0.5)
    add_shape_with_fill(slide, Inches(0.8), Inches(1.4), Inches(5.8), Pt(4), BLUE_PRIMARY)

    add_textbox(slide, Inches(1.2), Inches(1.6), Inches(5), Inches(0.4),
                "Contact Information", font_size=20, bold=True, color=NAVY_DARK)

    contacts = [
        ("Phone", "+91 93076 73391"),
        ("Email", "info@credsuvidha.com"),
        ("Website", "www.credsuvidha.com"),
        ("Social", "Facebook • Twitter/X • LinkedIn • Instagram"),
    ]

    for i, (label, value) in enumerate(contacts):
        y = Inches(2.2) + i * Inches(0.55)
        add_textbox(slide, Inches(1.4), y, Inches(1.5), Inches(0.3),
                    f"{label}:", font_size=13, bold=True, color=BLUE_PRIMARY)
        add_textbox(slide, Inches(3.0), y, Inches(3.5), Inches(0.3),
                    value, font_size=13, color=TEXT_DARK)

    # Compliance Card
    comp_card = add_shape_with_fill(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(3.5), WHITE)
    comp_card.line.color.rgb = BORDER; comp_card.line.width = Pt(0.5)
    add_shape_with_fill(slide, Inches(7.0), Inches(1.4), Inches(5.5), Pt(4), EMERALD)

    add_textbox(slide, Inches(7.4), Inches(1.6), Inches(5), Inches(0.4),
                "Compliance & Certifications", font_size=20, bold=True, color=NAVY_DARK)

    badges = [
        ("RBI Regulated Partners", "All banking partners regulated by Reserve Bank of India"),
        ("IRDAI Registered", "Insurance products through IRDAI registered entities"),
        ("ISO 27001 Compliant", "Information security management system compliance"),
    ]

    for i, (badge, desc) in enumerate(badges):
        y = Inches(2.2) + i * Inches(0.75)

        check = add_shape_with_fill(slide, Inches(7.4), y + Pt(2), Inches(0.3), Inches(0.3),
                                    EMERALD, MSO_SHAPE.OVAL)
        add_textbox(slide, Inches(7.4), y + Pt(2), Inches(0.3), Inches(0.3),
                    "✓", font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(7.9), y, Inches(4.3), Inches(0.3),
                    badge, font_size=13, bold=True, color=NAVY_DARK)
        add_textbox(slide, Inches(7.9), y + Inches(0.3), Inches(4.3), Inches(0.3),
                    desc, font_size=10, color=TEXT_GRAY)

    # Disclaimer
    disclaimer_box = add_shape_with_fill(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.8), BLUE_TINT)
    disclaimer_box.line.color.rgb = BLUE_PALE; disclaimer_box.line.width = Pt(0.5)

    add_textbox(slide, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.6),
                "Disclaimer: CredSuvidha acts as a facilitator. All financial products are subject to "
                "respective institution policies. Interest rates may vary based on individual eligibility and market conditions.",
                font_size=10, italic=True, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    notes = slide.notes_slide.notes_text_frame
    notes.text = ("[KEY POINT]: Contact us at +91 93076 73391 or info@credsuvidha.com.\n"
                  "[DATA]: We maintain RBI, IRDAI, and ISO 27001 compliance across all partners.\n"
                  "[TRANSITION]: Finally, let's review the technical specifications.")


def build_technical_slide(prs):
    """Slide 10: Technical Specifications"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_ALT

    add_accent_bar(slide)
    add_logo_header(slide)
    add_footer_bar(slide)
    add_slide_number(slide, 10)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                "Technical Specifications", font_size=32, bold=True, color=NAVY_DARK)

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4),
                "Technology stack and dependencies powering the CredSuvidha website.",
                font_size=13, color=TEXT_GRAY)

    specs = [
        ("CSS Framework", "Tailwind CSS via CDN\n(cdn.tailwindcss.com) with custom\nbrand theme configuration", BLUE_PRIMARY),
        ("Typography", "Google Fonts: Inter (300-900) +\nPlayfair Display (700-800)\nPreconnected for performance", LOGO_GOLD),
        ("JavaScript", "Vanilla JS — scroll reveal,\ncounter animation, mobile menu,\nsmooth scroll, form handling", EMERALD),
        ("Architecture", "Single-page HTML application\nNo build step required\nAll-in-one file deployment", ACCENT_PRI),
        ("Hosting", "Netlify (static hosting)\nCustom domain: www.credsuvidha.com\nGit-based deployment", NAVY),
        ("Version Control", "Git repository\nPushed to remote\nContinuous deployment via Netlify", BLUE_BRIGHT),
    ]

    card_w = Inches(3.7)
    card_h = Inches(2.1)

    for i, (title, desc, color) in enumerate(specs):
        col = i % 3
        row = i // 3
        x = Inches(0.8) + col * (card_w + Inches(0.25))
        y = Inches(1.65) + row * (card_h + Inches(0.25))

        card = add_shape_with_fill(slide, x, y, card_w, card_h, WHITE)
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5)

        # Left accent bar
        add_shape_with_fill(slide, x, y, Pt(5), card_h, color)

        add_textbox(slide, x + Inches(0.3), y + Inches(0.2), card_w - Inches(0.5), Inches(0.35),
                    title, font_size=15, bold=True, color=NAVY_DARK)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.6), card_w - Inches(0.5), Inches(1.3),
                    desc, font_size=12, color=TEXT_GRAY)

    # Asset inventory note
    assets_box = add_shape_with_fill(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.6), BLUE_TINT)
    assets_box.line.color.rgb = BLUE_PALE; assets_box.line.width = Pt(0.5)

    add_textbox(slide, Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.5),
                "Assets: logo.png • assets/images/logo.png • assets/brandkit/brand-guidelines.html • "
                "assets/brandkit/brand-tokens.json • assets/brandkit/CredSuvidha-BrandKit.pdf • "
                "assets/brandkit/CredSuvidha-BrandKit.pptx",
                font_size=9, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    notes = slide.notes_slide.notes_text_frame
    notes.text = ("[KEY POINT]: The website is a single HTML file deployed on Netlify with no build step.\n"
                  "[DATA]: Tailwind CSS + Google Fonts + Vanilla JS. All resources are CDN-based.\n"
                  "[TRANSITION]: Let's summarize what we've covered in this brand kit.")


def build_summary_slide(prs):
    """Slide 11: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BLUE_TINT

    add_accent_bar(slide, color=NAVY_DARK, height=Pt(5))
    add_logo_header(slide)
    add_footer_bar(slide)
    add_slide_number(slide, 11)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                "Key Takeaways", font_size=32, bold=True, color=NAVY_DARK)

    takeaways = [
        ("Logo", "Interlocking C&S with growth arrow — navy + gold on light backgrounds", BLUE_PRIMARY),
        ("Colors", "Primary #1A6EF5  •  Navy #142857  •  Accent #F97316  •  Emerald #10B981", NAVY),
        ("Typography", "Inter (sans-serif) for everything  •  Playfair Display (serif) for display", LOGO_GOLD),
        ("Buttons", "4 variants: Primary, Secondary, Accent, Dark — all pill-shaped (100px radius)", EMERALD),
        ("Values", "Trust, Speed, Expertise, Simplicity — every decision guided by these pillars", ACCENT_PRI),
        ("Compliance", "RBI Regulated  •  IRDAI Registered  •  ISO 27001 Compliant", BLUE_PRIMARY),
        ("Tech Stack", "Tailwind CSS + Inter/Playfair fonts + Vanilla JS → Netlify deployment", NAVY),
    ]

    for i, (label, desc, color) in enumerate(takeaways):
        y = Inches(1.4) + i * Inches(0.72)

        card = add_shape_with_fill(slide, Inches(0.8), y, Inches(11.7), Inches(0.6), WHITE)
        card.line.color.rgb = BORDER; card.line.width = Pt(0.5)

        # Left color bar
        add_shape_with_fill(slide, Inches(0.8), y, Pt(5), Inches(0.6), color)

        add_textbox(slide, Inches(1.2), y + Pt(4), Inches(1.5), Inches(0.4),
                    f"●  {label}", font_size=13, bold=True, color=color)
        add_textbox(slide, Inches(2.8), y + Pt(4), Inches(9.5), Inches(0.4),
                    desc, font_size=12, color=TEXT_DARK)

    # Bottom note
    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.35),
                "This brand kit is a living document — update as CredSuvidha evolves. Consistency builds trust.",
                font_size=12, bold=True, italic=True, color=NAVY, alignment=PP_ALIGN.CENTER)

    notes = slide.notes_slide.notes_text_frame
    notes.text = ("[KEY POINT]: Recap of all brand kit elements covered.\n"
                  "[DATA]: 7 key areas — logo, colors, typography, buttons, values, compliance, tech.\n"
                  "[TRANSITION]: Thank you for reviewing the CredSuvidha brand kit.")


def build_thankyou_slide(prs):
    """Slide 12: Thank You / Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY_DARK

    # Gold + Blue stripe
    add_shape_with_fill(slide, Inches(0), Inches(3.2), Inches(6.667), Pt(3), BLUE_PRIMARY)
    add_shape_with_fill(slide, Inches(6.667), Inches(3.2), Inches(6.666), Pt(3), LOGO_GOLD)

    # Logo
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(4.2), Inches(0.5), width=Inches(5.0))

    # Thank you text
    add_textbox(slide, Inches(1), Inches(3.5), Inches(11.333), Inches(0.9),
                "Thank You", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Contact info
    add_textbox(slide, Inches(1), Inches(4.5), Inches(11.333), Inches(0.5),
                "+91 93076 73391  •  info@credsuvidha.com  •  www.credsuvidha.com",
                font_size=16, color=BLUE_LIGHT, alignment=PP_ALIGN.CENTER)

    # Tagline
    txBox = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(11.333), Inches(0.5))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = "Trusted Partner. "; set_font(r1, 20, bold=True, color=WHITE)
    r2 = p.add_run(); r2.text = "Swift Solutions."; set_font(r2, 20, bold=True, color=LOGO_GOLD)

    # Copyright
    add_textbox(slide, Inches(1), Inches(6.3), Inches(11.333), Inches(0.3),
                "© 2025 CredSuvidha. All rights reserved.",
                font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # Bottom accent
    add_shape_with_fill(slide, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3), BLUE_PRIMARY)

    notes = slide.notes_slide.notes_text_frame
    notes.text = ("[KEY POINT]: Thank the audience for their time.\n"
                  "[DATA]: Provide contact details for follow-up questions.\n"
                  "For brand kit updates or questions, reach out to info@credsuvidha.com.")


# ═══════════════════════════════════════════════
# MAIN BUILD
# ═══════════════════════════════════════════════

def build_presentation():
    """Build the complete CredSuvidha Brand Kit presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_title_slide(prs)          # Slide 1
    build_brand_overview(prs)       # Slide 2
    build_logo_slide(prs)           # Slide 3
    build_color_palette_slide(prs)  # Slide 4
    build_gradients_slide(prs)      # Slide 5
    build_typography_slide(prs)     # Slide 6
    build_buttons_ui_slide(prs)     # Slide 7
    build_brand_values_slide(prs)   # Slide 8
    build_contact_compliance_slide(prs)  # Slide 9
    build_technical_slide(prs)      # Slide 10
    build_summary_slide(prs)        # Slide 11
    build_thankyou_slide(prs)       # Slide 12

    prs.save(OUTPUT_PATH)
    file_size = os.path.getsize(OUTPUT_PATH)
    print(f"Presentation generated successfully!")
    print(f"Output: {OUTPUT_PATH}")
    print(f"Slides: 12")
    print(f"Size: {file_size / 1024:.1f} KB")


if __name__ == '__main__':
    build_presentation()
